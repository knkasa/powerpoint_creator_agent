import json
import boto3
from io import BytesIO
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt

s3_client = boto3.client('s3')
bedrock_client = boto3.client('bedrock-runtime')

def lambda_handler(event, context):
    try:
        body = json.loads(event['body'])
        prompt = body['prompt']

        font_size = 18
        num_layout = 8
        num_page = 3

        system_prompt = f'''次の問いに対してプレゼンテーション用のパワーポイントを{num_page}ページ作成します。記載内容を下のフォーマットに沿って教えてください。 問い＝「{prompt}」
<1ページ>
<title>タイトル</title>
<text>
insert subtopic
  - insert texts
insert subtopic
  - insert texts
</text>
'''

        payload = {
            "anthropic_version": "bedrock-2023-05-31",
            "temperature": 0.5,
            "max_tokens": 5000,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": system_prompt
                        }
                    ]
                }
            ]
        }

        print("checking 1")
        response_raw = bedrock_client.invoke_model(
                modelId="anthropic.claude-3-haiku-20240307-v1:0",  # you can lookup from clicking "Model Catalog" in Bedrock webpage.
                body=json.dumps(payload),
                contentType="application/json",
                accept="application/json",
            )
        print("checking 2")

        response = json.loads(response_raw['body'].read())['content'][0]['text']
        print(response)

        s3_object = s3_client.get_object(Bucket='test-ecs-s3', Key='template.pptx')
        pptx_data = s3_object["Body"].read()
        print("template.pptx downloaded.")

        prs = Presentation(BytesIO(pptx_data))
        print("template.pptx loaded.")
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])  #0=title layout
        slide.shapes.title.text = prompt
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)

        for res in response.split('ページ')[1:]:

            soup = BeautifulSoup(res, 'html.parser')
            title_tag = soup.find('title')
            text_tag = soup.find('text')

            slide = prs.slides.add_slide(prs.slide_layouts[8])

            # Topic title
            title = slide.shapes.title
            title.text = title_tag.string.strip()
            title.text_frame.paragraphs[0].font.size = Pt(font_size+6)
            title.text_frame.paragraphs[0].font.bold = True
            title.top = Inches(-4)
            title.left = Inches(0.2)

            # Topic content
            content = slide.shapes.placeholders[1]
            content.text = text_tag.string.strip()
            content.top = Inches(1.5)
            content.left = Inches(0.3)
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)

        xml_slides = prs.slides._sldIdLst
        for _ in range(14): # delete template slides.
            xml_slides.remove(xml_slides[0]) 

        output_file="presentation.pptx"
        tmp_dir = f"/tmp/{output_file}" #Needs to be in /tmp folder.
        prs.save(tmp_dir)

        s3_client.upload_file(tmp_dir, 'test-ecs-s3', output_file)

        return {
            'statusCode': 200,
            'body': json.dumps("Job done.")
            }
        
    except Exception as e:
        print(e)
        return {
            'statusCode': 600,
            'body': json.dumps(str(e))
            }

